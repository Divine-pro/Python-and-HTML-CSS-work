import collections
import collections.abc
from pptx import Presentation

#read sample  slide data 1 and 2
r=open("sample_slide1_input.txt","r")
r1=open("sample_slide2_input.txt","r")

data1 = r.read()
data2 = r1.read()

# Create a presentation object
prs = Presentation()

font_path = 'sample_font_file.ttf'

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
text_box = slide.shapes.add_textbox(0, 0, 0, 0)
tf = text_box.text_frame

paragraph = tf.add_paragraph()
paragraph.text = data1
paragraph.font.file = font_path

slide1 = prs.slides.add_slide(title_slide_layout)
text_box1 = slide1.shapes.add_textbox(0, 0, 0, 0)
tf1 = text_box1.text_frame

paragraph1 = tf1.add_paragraph()
paragraph1.text = data2
paragraph1.font.file = font_path



#save the presentation
prs.save('sample_presentation.pptx')


